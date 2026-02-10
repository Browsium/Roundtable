import os
import tempfile
from pathlib import Path
from typing import Optional, Dict, Any
import logging

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Process uploaded documents and extract text."""
    
    SUPPORTED_EXTENSIONS = {'.pdf', '.docx', '.pptx', '.txt', '.md'}
    
    @classmethod
    def process_file(cls, file_path: str) -> Dict[str, Any]:
        """Process a file and return extracted text with metadata."""
        path = Path(file_path)
        
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        extension = path.suffix.lower()
        
        if extension not in cls.SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file type: {extension}")
        
        # Extract metadata
        stat = path.stat()
        metadata = {
            "filename": path.name,
            "size_bytes": stat.st_size,
            "extension": extension,
            "created_date": stat.st_ctime,
            "modified_date": stat.st_mtime
        }
        
        # Extract text based on file type
        if extension == '.pdf':
            text = cls._extract_pdf(file_path)
        elif extension == '.docx':
            text = cls._extract_docx(file_path)
        elif extension == '.pptx':
            text = cls._extract_pptx(file_path)
        elif extension in {'.txt', '.md'}:
            text = cls._extract_text(file_path)
        else:
            raise ValueError(f"Unsupported file type: {extension}")
        
        return {
            "text": text,
            "metadata": metadata
        }
    
    @staticmethod
    def _extract_pdf(file_path: str) -> str:
        """Extract text from PDF file."""
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            text_parts = []
            
            for page in reader.pages:
                text_parts.append(page.extract_text())
            
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"PDF extraction error: {e}")
            raise Exception(f"Failed to extract PDF: {e}")
    
    @staticmethod
    def _extract_docx(file_path: str) -> str:
        """Extract text from DOCX file."""
        try:
            from docx import Document
            doc = Document(file_path)
            
            text_parts = []
            for para in doc.paragraphs:
                text_parts.append(para.text)
            
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"DOCX extraction error: {e}")
            raise Exception(f"Failed to extract DOCX: {e}")
    
    @staticmethod
    def _extract_pptx(file_path: str) -> str:
        """Extract text from PPTX file."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"PPTX extraction error: {e}")
            raise Exception(f"Failed to extract PPTX: {e}")
    
    @staticmethod
    def _extract_text(file_path: str) -> str:
        """Extract text from plain text file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # Try with different encoding
            with open(file_path, 'r', encoding='latin-1') as f:
                return f.read()
        except Exception as e:
            logger.error(f"Text extraction error: {e}")
            raise Exception(f"Failed to extract text: {e}")
    
    @staticmethod
    def get_file_version(file_path: str) -> Optional[str]:
        """Try to extract version information from file metadata."""
        # This is a placeholder - in production you might extract
        # version from document properties or filename
        path = Path(file_path)
        
        # Check for version in filename (e.g., "doc-v2.pdf" or "doc_1.5.txt")
        import re
        version_match = re.search(r'[vV](\d+)|_(\d+\.\d+)', path.stem)
        if version_match:
            return version_match.group(1) or version_match.group(2)
        
        return None